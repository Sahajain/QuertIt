import streamlit as st
from langchain.text_splitter import RecursiveCharacterTextSplitter
import mimetypes
import os
import pandas as pd
from typing import Dict, List, Any, Optional
import sqlparse
import sqlite3

def detect_file_type(file):
    """
    Detect file type using both extension and MIME-type analysis
    """
    file_extension = os.path.splitext(file.name)[1].lower().replace('.', '')
    mime_type, _ = mimetypes.guess_type(file.name)
    
    return file_extension, mime_type

def get_chunking_params(file_type: str) -> Dict[str, Any]:
    """
    Return chunking parameters optimized for different document types
    based on empirical testing as described in the research methodology.
    """
    chunking_params = {
        "default": {
            "chunk_size": 800,
            "chunk_overlap": 150,
            "separators": ["\n\n", "\n", " ", ""]
        },
        "pdf": {
            "chunk_size": 800,
            "chunk_overlap": 150,
            "separators": ["\n\n", "\n", " ", ""]
        },
        "docx": {
            "chunk_size": 600,
            "chunk_overlap": 100,
            "separators": ["\n\n", "\n", ". ", " ", ""]
        },
        "txt": {
            "chunk_size": 1000,
            "chunk_overlap": 200,
            "separators": ["\n\n", "\n", ". ", " ", ""]
        },
        "pptx": {
            "chunk_size": 400,
            "chunk_overlap": 50,
            "separators": ["\n\n", "\n", ". ", " ", ""]
        },
        "csv": {
            "chunk_size": 500,
            "chunk_overlap": 100,
            "separators": ["\n\n", "\n", ". ", " ", ""]
        },
        "xlsx": {
            "chunk_size": 500,
            "chunk_overlap": 100,
            "separators": ["\n\n", "\n", ". ", " ", ""]
        },
        "sql": {
            "chunk_size": 600,
            "chunk_overlap": 100,
            "separators": [";\n", "\n\n", "\n", " ", ""]
        },
        "db": {
            "chunk_size": 600,
            "chunk_overlap": 100,
            "separators": ["\n\n", "\n", " ", ""]
        }
    }
    
    return chunking_params.get(file_type, chunking_params["default"])

def extract_text_from_pdf(file):
    """
    Extract text from PDF files with enhanced metadata extraction
    """
    from PyPDF2 import PdfReader
    
    text = ""
    pdf_reader = PdfReader(file)
    metadata = pdf_reader.metadata
    page_texts = []
    
    for i, page in enumerate(pdf_reader.pages):
        page_text = page.extract_text()
        if page_text:
            page_texts.append({"page_num": i+1, "text": page_text})
            text += f"\n\n--- Page {i+1} ---\n\n{page_text}"
    
    return text, {
        "title": metadata.title if metadata and metadata.title else None,
        "author": metadata.author if metadata and metadata.author else None,
        "total_pages": len(pdf_reader.pages)
    }

def extract_text_from_docx(file):
    """
    Extract text from DOCX files with section awareness
    """
    from docx import Document
    
    doc = Document(file)
    text = ""
    sections = []
    current_section = "Header"
    
    for para in doc.paragraphs:
        if para.style.name.startswith('Heading'):
            current_section = para.text
            sections.append(current_section)
        
        text += f"{para.text}\n"
    
    metadata = {
        "sections": sections,
        "total_paragraphs": len(doc.paragraphs)
    }
    
    return text, metadata

def extract_text_from_pptx(file):
    """
    Extract text from PPTX files with slide-based preservation
    """
    from pptx import Presentation
    
    presentation = Presentation(file)
    text = ""
    slides_content = []
    
    for i, slide in enumerate(presentation.slides):
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text += shape.text + "\n"
        
        slides_content.append({"slide_num": i+1, "text": slide_text})
        text += f"\n\n--- Slide {i+1} ---\n\n{slide_text}"
    
    metadata = {
        "total_slides": len(presentation.slides)
    }
    
    return text, metadata

def extract_text_from_txt(file):
    """
    Extract text from TXT files
    """
    content = file.read().decode("utf-8")
    
    paragraphs = content.split("\n\n")
    metadata = {
        "paragraphs": len(paragraphs)
    }
    
    return content, metadata

def extract_text_from_csv(file):
    """
    Extract text from CSV files with table preservation
    """
    import pandas as pd
    
    df = pd.read_csv(file)
    text = df.to_string()
    
    metadata = {
        "columns": list(df.columns),
        "rows": len(df)
    }
    
    return text, metadata

def extract_text_from_excel(file):
    """
    Extract text from Excel files with sheet awareness
    """
    import pandas as pd
    
    excel_file = pd.ExcelFile(file)
    text = ""
    sheets_data = {}
    
    for sheet_name in excel_file.sheet_names:
        df = pd.read_excel(excel_file, sheet_name=sheet_name)
        sheet_text = df.to_string()
        sheets_data[sheet_name] = {"rows": len(df), "columns": len(df.columns)}
        text += f"\n\n--- Sheet: {sheet_name} ---\n\n{sheet_text}"
    
    metadata = {
        "sheets": excel_file.sheet_names,
        "sheet_data": sheets_data
    }
    
    return text, metadata

def extract_text_from_image(file):
    """
    Extract text from images using OCR
    """
    try:
        from PIL import Image
        import pytesseract
        
        img = Image.open(file)
        text = pytesseract.image_to_string(img)
        
        metadata = {
            "width": img.width,
            "height": img.height,
            "format": img.format
        }
        
        return text, metadata
    except ImportError:
        st.warning("OCR dependencies (pytesseract) not installed. Skipping image text extraction.")
        return "", {}

def extract_text_from_markdown(file):
    """
    Extract text from Markdown files
    """
    try:
        import markdown
        
        content = file.read().decode("utf-8")
        html = markdown.markdown(content)
        
        headers = []
        for line in content.split("\n"):
            if line.startswith("#"):
                headers.append(line.strip("# "))
        
        metadata = {
            "headers": headers
        }
        
        return content, metadata
    except ImportError:
        return file.read().decode("utf-8"), {}

def extract_text_from_sql(file):
    """
    Extract text and metadata from SQL files
    """
    try:
        content = file.read().decode("utf-8")
        parsed = sqlparse.parse(content)
        
        tables = {}
        current_table = None
        
        for statement in parsed:
            if statement.get_type() == 'CREATE':
                tokens = statement.tokens
                for token in tokens:
                    if token.is_keyword and token.value.upper() == 'TABLE':
                        table_name = tokens[tokens.index(token) + 2].value
                        current_table = table_name
                        tables[current_table] = {'columns': [], 'sample_data': []}
                        for t in tokens:
                            if t.__class__.__name__ == 'Parenthesis':
                                for col in t.value.strip('()').split(','):
                                    col = col.strip()
                                    if col:
                                        tables[current_table]['columns'].append(col.split()[0])
            elif statement.get_type() == 'INSERT' and current_table:
                values = []
                for token in statement.tokens:
                    if token.__class__.__name__ == 'Parenthesis':
                        values.extend(token.value.strip('()').split(','))
                if len(tables[current_table]['sample_data']) < 5:
                    tables[current_table]['sample_data'].append([v.strip().strip("'") for v in values])
        
        metadata = {
            "tables": list(tables.keys()),
            "schema": {table: info['columns'] for table, info in tables.items()},
            "sample_data": {table: info['sample_data'] for table, info in tables.items()}
        }
        
        return content, metadata
    except Exception as e:
        st.error(f"Error parsing SQL file {file.name}: {str(e)}")
        return "", {}

def extract_text_from_db(file):
    """
    Extract text and metadata from SQLite .db files
    """
    try:
        with open("temp.db", "wb") as f:
            f.write(file.read())
        
        conn = sqlite3.connect("temp.db")
        cursor = conn.cursor()
        
        cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
        tables = [row[0] for row in cursor.fetchall()]
        
        text = ""
        metadata = {"tables": tables, "schema": {}, "sample_data": {}}
        
        for table in tables:
            cursor.execute(f"PRAGMA table_info({table});")
            columns = [row[1] for row in cursor.fetchall()]
            metadata["schema"][table] = columns
            
            cursor.execute(f"SELECT * FROM {table} LIMIT 5;")
            sample_data = cursor.fetchall()
            metadata["sample_data"][table] = sample_data
            
            text += f"\n\n--- Table: {table} ---\nColumns: {', '.join(columns)}\n"
            text += f"Sample Data:\n{pd.DataFrame(sample_data, columns=columns).to_string()}\n"
        
        conn.close()
        os.remove("temp.db")
        
        return text, metadata
    except Exception as e:
        st.error(f"Error processing SQLite file {file.name}: {str(e)}")
        return "", {}

def extract_metadata(file, file_type, raw_text):
    """
    Extract additional metadata based on file content and type
    """
    metadata = {
        "filename": file.name,
        "file_type": file_type,
        "character_count": len(raw_text)
    }
    
    return metadata

def get_page_number(chunk, raw_text, file_type):
    """
    Attempt to determine the page number for a given chunk
    """
    if "--- Page " in chunk:
        for line in chunk.split("\n"):
            if line.startswith("--- Page "):
                try:
                    return int(line.replace("--- Page ", "").replace(" ---", ""))
                except:
                    pass
    
    return None

def get_section(chunk, raw_text, file_type):
    """
    Attempt to determine the section for a given chunk
    """
    lines = chunk.split("\n")
    for line in lines:
        if line.startswith('#') or line.startswith("--- "):
            return line.strip('# -')
    
    return None

def get_text_chunks_with_metadata(files):
    """
    Process each file separately, extract text, chunk it, and assign metadata with source information.
    """
    all_chunks = []
    
    for file in files:
        file_extension, mime_type = detect_file_type(file)
        st.write(f"Processing file: {file.name} ({file_extension}, {mime_type})")
        
        try:
            if file_extension == "pdf":
                raw_text, file_metadata = extract_text_from_pdf(file)
            elif file_extension == "docx":
                raw_text, file_metadata = extract_text_from_docx(file)
            elif file_extension == "pptx":
                raw_text, file_metadata = extract_text_from_pptx(file)
            elif file_extension == "txt":
                raw_text, file_metadata = extract_text_from_txt(file)
            elif file_extension == "csv":
                raw_text, file_metadata = extract_text_from_csv(file)
            elif file_extension in ["xlsx", "xls"]:
                raw_text, file_metadata = extract_text_from_excel(file)
            elif file_extension in ["jpg", "jpeg", "png"]:
                raw_text, file_metadata = extract_text_from_image(file)
            elif file_extension in ["md", "html"]:
                raw_text, file_metadata = extract_text_from_markdown(file)
            elif file_extension == "sql":
                raw_text, file_metadata = extract_text_from_sql(file)
            elif file_extension == "db":
                raw_text, file_metadata = extract_text_from_db(file)
            else:
                st.warning(f"Unsupported file type: {file.name}")
                continue
                
            st.write(f"Extracted text length: {len(raw_text)}")
            
            chunking_params = get_chunking_params(file_extension)
            
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=chunking_params["chunk_size"],
                chunk_overlap=chunking_params["chunk_overlap"],
                separators=chunking_params["separators"],
                length_function=len
            )
            
            metadata = extract_metadata(file, file_extension, raw_text)
            metadata.update(file_metadata)
            
            chunks = text_splitter.split_text(raw_text)
            st.write(f"Number of chunks created: {len(chunks)}")
            
            for i, chunk in enumerate(chunks):
                page_num = get_page_number(chunk, raw_text, file_extension)
                section = get_section(chunk, raw_text, file_extension)
                
                chunk_metadata = {
                    "source": file.name,
                    "chunk_id": i,
                    "page": page_num,
                    "section": section,
                    **metadata
                }
                
                chunk_metadata = {k: v for k, v in chunk_metadata.items() if v is not None}
                
                all_chunks.append({
                    "page_content": chunk,
                    "metadata": chunk_metadata
                })
        
        except Exception as e:
            st.error(f"Error processing {file.name}: {str(e)}")
            continue
    
    return all_chunks